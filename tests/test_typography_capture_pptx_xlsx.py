# SPDX-License-Identifier: MIT
"""Cluster A PR-1: cross-format A1 CAPTURE (pptx + xlsx fill the docx shape).

The capture engine (``common.typography.capture_appearance`` /
``capture_palette_facts``) is format-neutral: each adapter (docx/pptx/xlsx) yields a
``RunFacts`` view and the engine records the SAME profile shape - the document body
font/size (``theme.fonts.body``), the body color (``theme.text.body``), the per-role
``appearance``, and the brand ``theme.palette``. These tests are the pptx/xlsx peers
of the docx ``CaptureTest`` / ``CaptureSizeColorTest`` in ``test_typography.py``.

They are UNIVERSAL by construction: every template here carries NON-default theme
fonts (Montserrat/Inter, never Calibri/Cambria) and NON-default brand slots, and the
captured run values are NON-default faces/sizes/hexes - so a capture that only worked
on the stock python-pptx / openpyxl template would fail. The dominant-value scoring,
the half-point unit (pptx sz/2 via Pt, xlsx round(sz*2)), and the THEME_SLOTS token
normalization are all proven against values that are not the library defaults.

A final SHAPE-PARITY test asserts the pptx/xlsx captured profile keys are exactly the
docx schema shape and that ``schema._validate_palette`` accepts every captured palette.
"""

from __future__ import annotations

import sys
import tempfile
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

import unittest

from openpyxl import Workbook
from openpyxl.styles import Color, Font
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

from brandkit.common import color as colorutil
from brandkit.common import typography as ct
from brandkit.formats.docx import typography as docx_typography
from brandkit.formats.pptx import typography as pptx_typography
from brandkit.formats.xlsx import typography as xlsx_typography
from brandkit.profile import schema

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Non-default theme injection (proves UNIVERSAL, not tuned to the lib defaults)
# ---------------------------------------------------------------------------
def _branded_theme_xml(major: str, minor: str, accent1_hex: str) -> bytes:
    """A minimal but valid DrawingML theme with NON-default fonts and an accent slot.

    The 12 clrScheme slots are all present (Office requires the full scheme); only
    ``accent1`` carries the test's brand hex, the rest are inert placeholders. The
    major/minor latin faces are the test's brand faces (Montserrat/Inter)."""
    slots = {
        "dk1": "000000",
        "lt1": "FFFFFF",
        "dk2": "1F1F1F",
        "lt2": "EEEEEE",
        "accent1": accent1_hex,
        "accent2": "111111",
        "accent3": "222222",
        "accent4": "333333",
        "accent5": "444444",
        "accent6": "555555",
        "hlink": "0000FF",
        "folHlink": "800080",
    }
    clr = "".join(
        f'<a:{slot}><a:srgbClr val="{hexv}"/></a:{slot}>'
        for slot, hexv in slots.items()
    )
    return (
        f'<a:theme xmlns:a="{_A_NS}" name="Brand">'
        f"<a:themeElements>"
        f'<a:clrScheme name="Brand">{clr}</a:clrScheme>'
        f'<a:fontScheme name="Brand">'
        f'<a:majorFont><a:latin typeface="{major}"/></a:majorFont>'
        f'<a:minorFont><a:latin typeface="{minor}"/></a:minorFont>'
        f"</a:fontScheme>"
        f'<a:fmtScheme name="Brand">'
        f'<a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        f'<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        f'<a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst>'
        f'<a:lnStyleLst><a:ln><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        f'<a:ln><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        f'<a:ln><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst>'
        f"<a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle>"
        f"<a:effectStyle><a:effectLst/></a:effectStyle>"
        f"<a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst>"
        f'<a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        f'<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        f'<a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst>'
        f"</a:fmtScheme>"
        f"</a:themeElements></a:theme>"
    ).encode("utf-8")


def _rewrite_part(src: Path, dst: Path, part: str, data: bytes) -> None:
    """Copy ``src`` zip to ``dst`` with ``part`` replaced by ``data`` (theme inject)."""
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            payload = data if item.filename == part else zin.read(item.filename)
            zout.writestr(item, payload)


# ---------------------------------------------------------------------------
# pptx capture (peers of docx CaptureTest / CaptureSizeColorTest)
# ---------------------------------------------------------------------------
def _branded_deck(
    path: Path, *, runs: int = 5, font="Montserrat", pt=14, hexc="C0392B"
):
    """A deck whose content body carries explicit NON-default run typography, with a
    NON-default brand theme (Montserrat/Inter, accent1 = brand hex) injected."""
    prs = Presentation()
    prs.slide_layouts[1].name = "BrandContent"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = next(ph for ph in slide.placeholders if ph.placeholder_format.idx != 0)
    tf = body.text_frame
    for i in range(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = "brand body line"
        run.font.name = font
        run.font.size = Pt(pt)
        run.font.color.rgb = RGBColor.from_string(hexc)
    raw = path.with_suffix(".raw.pptx")
    prs.save(raw)
    _rewrite_part(
        raw, path, "ppt/theme/theme1.xml", _branded_theme_xml(font, "Inter", hexc)
    )


def _theme_for(path: Path, slot_part: str) -> dict:
    """The seeded theme block for a deck/workbook, mirroring the extractor's _theme."""
    theme = {
        "colors": {
            slot: {"hex": hexv}
            for slot, hexv in colorutil.parse_theme_colors(
                zipfile.ZipFile(path).read(slot_part)
            ).items()
        },
        "palette_roles": {"primary": {"theme": "accent1"}, "text": {"theme": "dk1"}},
        "fonts": {
            "major": {"latin": None, "fallback": "Arial"},
            "minor": {"latin": None, "fallback": "Calibri"},
        },
        "embedded_fonts": [],
    }
    colorutil.seed_theme_palette(theme)
    return theme


class PptxCaptureTest(unittest.TestCase):
    def test_dominant_body_font_size_color_captured_into_same_shape(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _branded_deck(deck, font="Montserrat", pt=14, hexc="C0392B")
            prs = Presentation(deck)
            roles = {"_index": []}
            theme = _theme_for(deck, "ppt/theme/theme1.xml")
            ct.capture_appearance(
                pptx_typography.iter_run_facts(prs),
                roles,
                theme,
                role_style_key=lambda _e: None,
            )
            ct.capture_palette_facts(pptx_typography.iter_run_facts(prs), roles, theme)
            body = theme["fonts"]["body"]
            self.assertEqual(body["latin"], "Montserrat")  # NON-default face
            self.assertEqual(body["size_hp"], 28)  # 14pt * 2 (Pt unit)
            self.assertGreaterEqual(body["confidence"], 0.6)
            self.assertEqual(
                theme["text"]["body"]["color"], {"kind": "hex", "hex": "C0392B"}
            )
            self.assertEqual(schema._validate_palette(theme["palette"]), [])

    def test_run_scheme_color_normalized_to_theme_slot(self):
        # A SCHEME (theme) run color is captured as a canonical clrScheme slot token,
        # not a docx WML token, and folds into theme.palette idempotently over the seed.
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            prs = Presentation()
            prs.slide_layouts[1].name = "BrandContent"
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            body = next(
                ph for ph in slide.placeholders if ph.placeholder_format.idx != 0
            )
            tf = body.text_frame
            for i in range(5):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = para.add_run()
                run.text = "themed"
                run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
            raw = Path(t) / "deck.raw.pptx"
            prs.save(raw)
            _rewrite_part(
                raw,
                deck,
                "ppt/theme/theme1.xml",
                _branded_theme_xml("Montserrat", "Inter", "C0392B"),
            )
            prs2 = Presentation(deck)
            roles = {"_index": []}
            theme = _theme_for(deck, "ppt/theme/theme1.xml")
            ct.capture_appearance(
                pptx_typography.iter_run_facts(prs2),
                roles,
                theme,
                role_style_key=lambda _e: None,
            )
            ct.capture_palette_facts(pptx_typography.iter_run_facts(prs2), roles, theme)
            self.assertEqual(
                theme["text"]["body"]["color"], {"kind": "theme", "theme": "accent2"}
            )
            self.assertIn("accent2", theme["palette"])
            self.assertEqual(schema._validate_palette(theme["palette"]), [])

    def test_no_explicit_run_typography_captures_nothing(self):
        # A deck whose body inherits its typography (no explicit run face/size/color)
        # captures nothing - byte-identical no-op territory.
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            prs = Presentation()
            prs.slide_layouts[1].name = "BrandContent"
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            body = next(
                ph for ph in slide.placeholders if ph.placeholder_format.idx != 0
            )
            tf = body.text_frame
            for i in range(5):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.add_run().text = "inherits everything"
            raw = Path(t) / "deck.raw.pptx"
            prs.save(raw)
            _rewrite_part(
                raw,
                deck,
                "ppt/theme/theme1.xml",
                _branded_theme_xml("Montserrat", "Inter", "C0392B"),
            )
            prs2 = Presentation(deck)
            roles = {"_index": []}
            theme = _theme_for(deck, "ppt/theme/theme1.xml")
            ct.capture_appearance(
                pptx_typography.iter_run_facts(prs2),
                roles,
                theme,
                role_style_key=lambda _e: None,
            )
            self.assertNotIn("body", theme["fonts"])
            self.assertNotIn("text", theme)


# ---------------------------------------------------------------------------
# xlsx capture (peers of docx CaptureTest / CaptureSizeColorTest)
# ---------------------------------------------------------------------------
def _branded_workbook(path: Path, *, rows=5, font="Inter", pt=12, hexc="1F4E79"):
    """A workbook whose cells carry explicit NON-default font typography, with a
    NON-default brand theme injected (Montserrat/Inter, accent1 brand hex)."""
    wb = Workbook()
    ws = wb.active
    for i in range(rows):
        cell = ws.cell(row=i + 1, column=1, value="brand cell text")
        cell.font = Font(name=font, size=pt, color=Color(rgb="FF" + hexc))
    raw = path.with_suffix(".raw.xlsx")
    wb.save(raw)
    _rewrite_part(
        raw,
        path,
        "xl/theme/theme1.xml",
        _branded_theme_xml("Montserrat", font, "C0392B"),
    )


class XlsxCaptureTest(unittest.TestCase):
    def test_dominant_cell_font_size_color_captured_into_same_shape(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _branded_workbook(book, font="Inter", pt=12, hexc="1F4E79")
            from openpyxl import load_workbook

            wb = load_workbook(book)
            roles = {"_index": []}
            theme = _theme_for(book, "xl/theme/theme1.xml")
            ct.capture_appearance(xlsx_typography.iter_run_facts(wb), roles, theme)
            ct.capture_palette_facts(xlsx_typography.iter_run_facts(wb), roles, theme)
            body = theme["fonts"]["body"]
            self.assertEqual(body["latin"], "Inter")  # NON-default face
            self.assertEqual(body["size_hp"], 24)  # round(12 * 2)
            self.assertGreaterEqual(body["confidence"], 0.6)
            self.assertEqual(
                theme["text"]["body"]["color"], {"kind": "hex", "hex": "1F4E79"}
            )  # ARGB alpha stripped
            self.assertEqual(schema._validate_palette(theme["palette"]), [])

    def test_theme_index_color_maps_to_base_slot_dropping_tint(self):
        # An openpyxl theme color is an integer index into the Excel cell theme-color
        # order; capture maps it via _XLSX_THEME_INDEX to the base slot, dropping .tint.
        # index 4 (accent1) is in the unchanged tail, so it coincides with THEME_SLOTS.
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            wb = Workbook()
            ws = wb.active
            for i in range(5):
                cell = ws.cell(row=i + 1, column=1, value="themed")
                # index 4 -> accent1; a non-zero tint must be dropped (base slot kept)
                cell.font = Font(color=Color(theme=4, tint=0.4))
            raw = Path(t) / "book.raw.xlsx"
            wb.save(raw)
            _rewrite_part(
                raw,
                book,
                "xl/theme/theme1.xml",
                _branded_theme_xml("Montserrat", "Inter", "C0392B"),
            )
            from openpyxl import load_workbook

            wb2 = load_workbook(book)
            roles = {"_index": []}
            theme = _theme_for(book, "xl/theme/theme1.xml")
            ct.capture_appearance(xlsx_typography.iter_run_facts(wb2), roles, theme)
            ct.capture_palette_facts(xlsx_typography.iter_run_facts(wb2), roles, theme)
            self.assertEqual(
                theme["text"]["body"]["color"], {"kind": "theme", "theme": "accent1"}
            )
            self.assertIn("accent1", theme["palette"])
            self.assertEqual(schema._validate_palette(theme["palette"]), [])

    def test_theme_index_swap_zone_maps_to_correct_slot(self):
        # Regression: Excel's cell <color theme="N"> index swaps the first two
        # dark/light pairs vs clrScheme document order. theme=1 is Excel's DEFAULT
        # text color and must map to dk1 (dark), NOT lt1 (light) - else body text is
        # branded white/invisible. A raw THEME_SLOTS[1] lookup would (wrongly) yield lt1.
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            wb = Workbook()
            ws = wb.active
            for i in range(5):
                cell = ws.cell(row=i + 1, column=1, value="body text")
                cell.font = Font(color=Color(theme=1))  # Excel default Text 1
            raw = Path(t) / "book.raw.xlsx"
            wb.save(raw)
            _rewrite_part(
                raw,
                book,
                "xl/theme/theme1.xml",
                _branded_theme_xml("Montserrat", "Inter", "C0392B"),
            )
            from openpyxl import load_workbook

            wb2 = load_workbook(book)
            roles = {"_index": []}
            theme = _theme_for(book, "xl/theme/theme1.xml")
            ct.capture_appearance(xlsx_typography.iter_run_facts(wb2), roles, theme)
            self.assertEqual(
                theme["text"]["body"]["color"], {"kind": "theme", "theme": "dk1"}
            )

    def test_no_explicit_cell_typography_captures_nothing(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            wb = Workbook()
            ws = wb.active
            for i in range(5):
                ws.cell(row=i + 1, column=1, value="inherits everything")
            raw = Path(t) / "book.raw.xlsx"
            wb.save(raw)
            _rewrite_part(
                raw,
                book,
                "xl/theme/theme1.xml",
                _branded_theme_xml("Montserrat", "Inter", "C0392B"),
            )
            from openpyxl import load_workbook

            wb2 = load_workbook(book)
            roles = {"_index": []}
            theme = _theme_for(book, "xl/theme/theme1.xml")
            ct.capture_appearance(xlsx_typography.iter_run_facts(wb2), roles, theme)
            self.assertNotIn("body", theme["fonts"])
            self.assertNotIn("text", theme)


# ---------------------------------------------------------------------------
# Shape parity: pptx/xlsx fill the EXACT docx-captured shape
# ---------------------------------------------------------------------------
class CaptureShapeParityTest(unittest.TestCase):
    """The keys pptx/xlsx capture into ``theme.fonts.body`` / ``theme.text.body`` and
    the palette-entry shape are EXACTLY the docx shape - one engine, one shape."""

    def _docx_captured_theme(self):
        from docx import Document
        from docx.shared import Pt as DPt
        from docx.shared import RGBColor as DRGB

        doc = Document()
        for _ in range(5):
            run = doc.add_paragraph().add_run("body text")
            run.font.name = "Montserrat"
            run.font.size = DPt(14)
            run.font.color.rgb = DRGB.from_string("C0392B")
        theme = {"colors": {"accent1": {"hex": "C0392B"}}, "fonts": {}}
        docx_typography.capture_fonts(doc, {"_index": []}, theme)
        docx_typography.capture_palette(doc, {"_index": []}, theme)
        return theme

    def _pptx_captured_theme(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _branded_deck(deck, font="Montserrat", pt=14, hexc="C0392B")
            prs = Presentation(deck)
            theme = _theme_for(deck, "ppt/theme/theme1.xml")
            ct.capture_appearance(
                pptx_typography.iter_run_facts(prs),
                {"_index": []},
                theme,
                role_style_key=lambda _e: None,
            )
            ct.capture_palette_facts(
                pptx_typography.iter_run_facts(prs), {"_index": []}, theme
            )
            return theme

    def _xlsx_captured_theme(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _branded_workbook(book, font="Inter", pt=12, hexc="1F4E79")
            from openpyxl import load_workbook

            wb = load_workbook(book)
            theme = _theme_for(book, "xl/theme/theme1.xml")
            ct.capture_appearance(
                xlsx_typography.iter_run_facts(wb), {"_index": []}, theme
            )
            ct.capture_palette_facts(
                xlsx_typography.iter_run_facts(wb), {"_index": []}, theme
            )
            return theme

    def test_body_font_keys_match_across_formats(self):
        docx_body = self._docx_captured_theme()["fonts"]["body"]
        for theme in (self._pptx_captured_theme(), self._xlsx_captured_theme()):
            self.assertEqual(set(theme["fonts"]["body"]), set(docx_body))
            self.assertEqual(set(theme["text"]["body"]), {"color", "color_confidence"})

    def test_palette_entry_shape_matches_across_formats(self):
        docx_keyset = {
            "ref",
            "provenance",
            "frequency",
            "name",
            "purpose",
            "use_when",
        }
        for theme in (self._pptx_captured_theme(), self._xlsx_captured_theme()):
            self.assertTrue(theme["palette"])
            for entry in theme["palette"].values():
                self.assertEqual(set(entry), docx_keyset)
            # The validator (schema 1.2.0, additive) accepts every captured palette.
            self.assertEqual(schema._validate_palette(theme["palette"]), [])


if __name__ == "__main__":
    unittest.main()
