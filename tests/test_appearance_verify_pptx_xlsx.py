# SPDX-License-Identifier: MIT
"""Cluster A PR-2: format-neutral A3 VERIFY for pptx + xlsx.

``check_appearance_targets`` now dispatches via ``_SHELL_APPEARANCE_COLLECTORS`` and
feeds the per-kind :class:`ShellAppearanceFacts` into the UNCHANGED membership loops.
These tests are the pptx/xlsx peers of the docx ``AppearanceSizeColorCheckTest`` /
``AppearanceTargetsCheckTest`` in ``test_typography.py``: an applied font/size/color
the shell PROVES it carries verifies clean (``[]``); one the shell does NOT carry
fails closed (ERROR).

UNIVERSAL by construction: every synthetic shell carries NON-default brand theme
fonts (Montserrat/Inter) and a NON-default brand accent hex, and the asserted values
are not the library defaults - a check tuned to the stock python-pptx / openpyxl
template would not pass here. The highest-risk axis is the UNIT conversion: an 18pt
pptx run is ``a:rPr@sz="1800"`` (centipoints), which must verify as 36 half-points
(``1800 / 50``); an xlsx ``font.sz=18`` must verify as ``round(18 * 2) = 36``. A
mismatch would make a correctly-applied size spuriously fail closed.
"""

from __future__ import annotations

import sys
import tempfile
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

import unittest

from openpyxl import Workbook
from openpyxl.styles import Color, Font, NamedStyle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from brandkit.profile import schema
from brandkit.qa import checks_deterministic as cd

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Non-default brand theme injection (proves UNIVERSAL, not tuned to lib defaults)
# ---------------------------------------------------------------------------
def _branded_theme_xml(major: str, minor: str, accent1_hex: str) -> bytes:
    slots = {
        "dk1": "000000",
        "lt1": "FFFFFF",
        "dk2": "1F1F1F",
        "lt2": "EEEEEE",
        "accent1": accent1_hex,
        "accent2": "111111",
        "accent3": "222222",
        "accent4": "333333",
        "accent5": "444444",
        "accent6": "555555",
        "hlink": "0000FF",
        "folHlink": "800080",
    }
    clr = "".join(
        f'<a:{slot}><a:srgbClr val="{hexv}"/></a:{slot}>'
        for slot, hexv in slots.items()
    )
    return (
        f'<a:theme xmlns:a="{_A_NS}" name="Brand"><a:themeElements>'
        f'<a:clrScheme name="Brand">{clr}</a:clrScheme>'
        f'<a:fontScheme name="Brand">'
        f'<a:majorFont><a:latin typeface="{major}"/></a:majorFont>'
        f'<a:minorFont><a:latin typeface="{minor}"/></a:minorFont>'
        f"</a:fontScheme></a:themeElements></a:theme>"
    ).encode("utf-8")


def _rewrite_part(src: Path, dst: Path, part: str, data: bytes) -> None:
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            payload = data if item.filename == part else zin.read(item.filename)
            zout.writestr(item, payload)


def _errors(findings):
    return [
        f
        for f in findings
        if f.check == "appearance_targets_exist"
        and f.severity == schema.Severity.ERROR.value
    ]


# ---------------------------------------------------------------------------
# Synthetic shell builders
# ---------------------------------------------------------------------------
def _deck(path: Path, *, font="Montserrat", pt=18, hexc="C0392B"):
    """A deck whose body run carries explicit brand typography, with a NON-default
    brand theme injected. An 18pt run is ``a:rPr@sz="1800"`` (the unit regression)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = next(ph for ph in slide.placeholders if ph.placeholder_format.idx != 0)
    run = body.text_frame.paragraphs[0].add_run()
    run.text = "brand body line"
    run.font.name = font
    run.font.size = Pt(pt)
    run.font.color.rgb = RGBColor.from_string(hexc)
    raw = path.with_suffix(".raw.pptx")
    prs.save(raw)
    _rewrite_part(
        raw, path, "ppt/theme/theme1.xml", _branded_theme_xml(font, "Inter", hexc)
    )


def _workbook(path: Path, *, font="Inter", pt=18, hexc="1F4E79", named=False):
    """A workbook whose cell carries explicit brand typography, with a NON-default
    brand theme injected. ``named`` registers a NamedStyle whose font widens the
    allow-set (the xlsx peer of the docx fontTable widening)."""
    wb = Workbook()
    ws = wb.active
    cell = ws.cell(row=1, column=1, value="brand cell text")
    cell.font = Font(name=font, size=pt, color=Color(rgb="FF" + hexc))
    if named:
        style = NamedStyle(name="BrandStyle")
        style.font = Font(name="NamedOnlyFace", size=pt)
        wb.add_named_style(style)
    raw = path.with_suffix(".raw.xlsx")
    wb.save(raw)
    _rewrite_part(
        raw,
        path,
        "xl/theme/theme1.xml",
        _branded_theme_xml("Montserrat", font, "C0392B"),
    )


def _profile(kind: str, *, theme=None, roles=None) -> dict:
    prof = schema.build_envelope(kind, {"name": "t"})
    prof["surface"] = {kind: {}}
    prof["roles"] = roles or {"_index": []}
    if theme is not None:
        prof["theme"] = theme
    return prof


# ---------------------------------------------------------------------------
# pptx verify (peer of docx AppearanceSizeColorCheckTest)
# ---------------------------------------------------------------------------
class PptxVerifyTest(unittest.TestCase):
    def test_applied_font_size_color_present_in_shell_is_accepted(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck, font="Montserrat", pt=18, hexc="C0392B")
            prof = _profile(
                "pptx",
                theme={
                    "colors": {},
                    "fonts": {"body": {"latin": "Montserrat", "size_hp": 36}},
                    "text": {"body": {"color": {"kind": "hex", "hex": "C0392B"}}},
                },
            )
            self.assertEqual(cd.check_appearance_targets(deck, prof), [])

    def test_eighteen_pt_run_verifies_as_thirty_six_half_points(self):
        # The unit regression: a:rPr@sz="1800" (centipoints) -> 36 half-points (/50).
        # An applied 36 must be ACCEPTED; an applied 35/37 (one half-point off) ERROR.
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck, font="Montserrat", pt=18, hexc="C0392B")
            facts = cd._pptx_collect_appearance_facts(
                deck, {"kind": "pptx", "theme": {}}
            )
            self.assertIn(36, facts.sizes)
            self.assertNotIn(1800, facts.sizes)  # never the raw centipoint value
            ok = _profile(
                "pptx", theme={"colors": {}, "fonts": {"body": {"size_hp": 36}}}
            )
            self.assertEqual(cd.check_appearance_targets(deck, ok), [])
            bad = _profile(
                "pptx", theme={"colors": {}, "fonts": {"body": {"size_hp": 35}}}
            )
            self.assertTrue(_errors(cd.check_appearance_targets(deck, bad)))

    def test_off_template_size_is_error(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck, pt=18)
            prof = _profile(
                "pptx", theme={"colors": {}, "fonts": {"body": {"size_hp": 99}}}
            )
            self.assertTrue(_errors(cd.check_appearance_targets(deck, prof)))

    def test_off_shell_font_is_error(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck, font="Montserrat")
            prof = _profile(
                "pptx", theme={"colors": {}, "fonts": {"body": {"latin": "ZZZ Bogus"}}}
            )
            self.assertTrue(_errors(cd.check_appearance_targets(deck, prof)))

    def test_off_palette_hex_color_is_error(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck, hexc="C0392B")
            prof = _profile(
                "pptx",
                theme={
                    "colors": {},
                    "fonts": {},
                    "text": {"body": {"color": {"kind": "hex", "hex": "ABCDEF"}}},
                },
            )
            self.assertTrue(_errors(cd.check_appearance_targets(deck, prof)))

    def test_palette_theme_token_color_is_accepted(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck, hexc="C0392B")
            prof = _profile(
                "pptx",
                theme={
                    "colors": {},
                    "fonts": {},
                    "text": {"body": {"color": {"kind": "theme", "theme": "accent1"}}},
                },
            )
            self.assertEqual(cd.check_appearance_targets(deck, prof), [])

    def test_empty_appearance_profile_has_no_finding(self):
        with tempfile.TemporaryDirectory() as t:
            deck = Path(t) / "deck.pptx"
            _deck(deck)
            prof = _profile("pptx", theme={"colors": {}, "fonts": {}})
            self.assertEqual(cd.check_appearance_targets(deck, prof), [])


# ---------------------------------------------------------------------------
# xlsx verify (peer of docx AppearanceSizeColorCheckTest)
# ---------------------------------------------------------------------------
class XlsxVerifyTest(unittest.TestCase):
    def test_applied_font_size_color_present_in_shell_is_accepted(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, font="Inter", pt=18, hexc="1F4E79")
            prof = _profile(
                "xlsx",
                theme={
                    "colors": {},
                    "fonts": {"body": {"latin": "Inter", "size_hp": 36}},
                    "text": {"body": {"color": {"kind": "hex", "hex": "1F4E79"}}},
                },
            )
            self.assertEqual(cd.check_appearance_targets(book, prof), [])

    def test_named_style_font_widens_the_allow_set(self):
        # A font that lives ONLY on a NamedStyle (no cell uses it) is still ALLOWED -
        # the xlsx peer of the docx fontTable/theme widening.
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, named=True)
            facts = cd._xlsx_collect_appearance_facts(
                book, {"kind": "xlsx", "theme": {}}
            )
            self.assertIn("NamedOnlyFace", facts.fonts)
            prof = _profile(
                "xlsx",
                theme={"colors": {}, "fonts": {"body": {"latin": "NamedOnlyFace"}}},
            )
            self.assertEqual(cd.check_appearance_targets(book, prof), [])

    def test_size_round_trips_through_half_points(self):
        # font.sz=18 -> round(18*2)=36 half-points. An applied 36 is accepted.
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, pt=18)
            facts = cd._xlsx_collect_appearance_facts(
                book, {"kind": "xlsx", "theme": {}}
            )
            self.assertIn(36, facts.sizes)
            ok = _profile(
                "xlsx", theme={"colors": {}, "fonts": {"body": {"size_hp": 36}}}
            )
            self.assertEqual(cd.check_appearance_targets(book, ok), [])

    def test_off_template_size_is_error(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, pt=18)
            prof = _profile(
                "xlsx", theme={"colors": {}, "fonts": {"body": {"size_hp": 99}}}
            )
            self.assertTrue(_errors(cd.check_appearance_targets(book, prof)))

    def test_off_shell_font_is_error(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, font="Inter")
            prof = _profile(
                "xlsx", theme={"colors": {}, "fonts": {"body": {"latin": "ZZZ Bogus"}}}
            )
            self.assertTrue(_errors(cd.check_appearance_targets(book, prof)))

    def test_run_provenance_hex_color_is_accepted(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, hexc="1F4E79")  # not a palette slot, but on a cell font
            prof = _profile(
                "xlsx",
                theme={
                    "colors": {},
                    "fonts": {},
                    "text": {"body": {"color": {"kind": "hex", "hex": "1F4E79"}}},
                },
            )
            self.assertEqual(cd.check_appearance_targets(book, prof), [])

    def test_off_template_hex_color_is_error(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book, hexc="1F4E79")
            prof = _profile(
                "xlsx",
                theme={
                    "colors": {},
                    "fonts": {},
                    "text": {"body": {"color": {"kind": "hex", "hex": "ABCDEF"}}},
                },
            )
            self.assertTrue(_errors(cd.check_appearance_targets(book, prof)))

    def test_empty_appearance_profile_has_no_finding(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book)
            prof = _profile("xlsx", theme={"colors": {}, "fonts": {}})
            self.assertEqual(cd.check_appearance_targets(book, prof), [])


# ---------------------------------------------------------------------------
# Fail-closed: a collector that cannot parse -> WARNING + every applied value ERROR
# ---------------------------------------------------------------------------
class FailClosedTest(unittest.TestCase):
    def test_unparseable_shell_warns_and_fails_applied_value_closed(self):
        with tempfile.TemporaryDirectory() as t:
            garbage = Path(t) / "broken.pptx"
            garbage.write_bytes(b"not a zip at all")
            prof = _profile(
                "pptx",
                theme={
                    "colors": {},
                    "fonts": {"body": {"latin": "Montserrat", "size_hp": 36}},
                    "text": {"body": {"color": {"kind": "hex", "hex": "C0392B"}}},
                },
            )
            findings = cd.check_appearance_targets(garbage, prof)
            warns = [
                f
                for f in findings
                if f.check == "appearance_targets_exist"
                and f.severity == schema.Severity.WARNING.value
            ]
            self.assertTrue(warns, [f.message for f in findings])
            # Empty fact sets -> the applied font/size/color all fail closed.
            self.assertTrue(_errors(findings))

    def test_unknown_kind_is_a_noop(self):
        with tempfile.TemporaryDirectory() as t:
            book = Path(t) / "book.xlsx"
            _workbook(book)
            prof = _profile(
                "xlsx", theme={"colors": {}, "fonts": {"body": {"latin": "X"}}}
            )
            prof["kind"] = "pdf"  # no registered collector
            self.assertEqual(cd.check_appearance_targets(book, prof), [])

    def test_none_shell_is_a_noop(self):
        prof = _profile("pptx", theme={"colors": {}, "fonts": {"body": {"latin": "X"}}})
        self.assertEqual(cd.check_appearance_targets(None, prof), [])


if __name__ == "__main__":
    unittest.main()
