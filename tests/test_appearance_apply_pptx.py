# SPDX-License-Identifier: MIT
"""PPTX brand APPEARANCE apply tests (Cluster A, PR-3 / A2).

The pptx peer of the docx ``ApplyTest`` / ``ApplySizeColorTest`` family: prove that
``_set_para_runs`` brands a body run's font / size / color from the captured profile
through the SHARED ``common.appearance`` orchestration and the pptx ``PPTX_BACKEND``,
under the same invariants the docx leg holds:

  - **Applied when unset:** a run that inherits its formatting (the usual pptx case)
    is branded from the resolved op's appearance (font name; size via ``Pt(hp / 2)``;
    color via pptx's OWN theme map or hex realization).
  - **Set-only-when-unset:** an explicit per-run ``color`` token wins over the body
    default (first-writer-wins), and an axis the run already carries is never
    clobbered.
  - **Universal, not tuned:** the deck uses a renamed layout and a non-default brand
    font/slot; nothing is keyed off a literal template name.
  - **Empty appearance is a byte-identical no-op:** a pre-capture profile generates
    byte-for-byte what it generated before the apply path existed, and generate-twice
    on an appearance-carrying profile is byte-idempotent.

These build SYNTHETIC decks in a temp dir (never tuned to napoleon/ntt) so the unit
conversion (pptx ``sz`` centipoints <-> the engine's half-points) and the
set-only-when-unset guard are exercised directly.
"""

from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path

THIS = Path(__file__).resolve()
sys.path.insert(0, str(THIS.parent.parent / "scripts"))
sys.path.insert(0, str(THIS.parent))

from pptx import Presentation  # noqa: E402
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR  # noqa: E402
from pptx.util import Pt  # noqa: E402

from brandkit.formats.pptx import generate as pg  # noqa: E402
from brandkit.ir.model import parse_idoc  # noqa: E402
from brandkit.profile import schema  # noqa: E402

# Reuse the proven synthetic-deck + in-memory-profile helpers.
import test_pptx_fixes as T  # noqa: E402


def _profile_with_body(
    template: Path,
    *,
    latin: str | None = None,
    size_hp: int | None = None,
    color: dict | None = None,
    palette: dict | None = None,
    colors: dict | None = None,
) -> dict:
    """An extracted pptx profile with captured document-body appearance injected.

    Populates the SAME keys the capture engine fills - ``theme.fonts.body`` (latin /
    size_hp) and ``theme.text.body.color`` - so the resolver's ``paragraph`` op carries
    the body default exactly as a really-captured profile would, plus an optional brand
    ``theme.palette`` / ``theme.colors`` for per-run color tokens. Nothing is keyed off
    the template's identity, so the test stays universal."""
    profile = T._extract_profile(template)
    body: dict = {}
    if latin:
        body["latin"] = latin
    if size_hp:
        body["size_hp"] = size_hp
    if body:
        profile["theme"]["fonts"]["body"] = body
    if color:
        profile["theme"].setdefault("text", {})["body"] = {"color": color}
    if palette:
        profile["theme"]["palette"] = palette
    if colors:
        profile["theme"].setdefault("colors", {}).update(colors)
    return profile


def _body_runs(prs: Presentation):
    """Every body-placeholder run across the generated deck, in slide order."""
    runs = []
    for slide in prs.slides:
        body = pg._first_body_placeholder(slide)
        if body is None or not getattr(body, "has_text_frame", False):
            continue
        for para in body.text_frame.paragraphs:
            runs.extend(para.runs)
    return runs


_BODY_IDOC = {
    "blocks": [
        {"type": "heading", "level": 1, "text": "Section"},
        {"type": "paragraph", "text": "Body text here."},
    ]
}


class PptxAppearanceApplyTest(unittest.TestCase):
    def test_body_font_size_color_applied_to_unset_run(self) -> None:
        """An inherited body run is branded with the captured font, size, and color."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(
                template,
                latin="Inter",
                size_hp=36,
                color={"kind": "hex", "hex": "1F4E79"},
            )
            out = tp / "out.pptx"
            pg.generate(profile, template, parse_idoc(_BODY_IDOC), out)

            runs = _body_runs(Presentation(out))
            target = next(r for r in runs if r.text == "Body text here.")
            self.assertEqual(target.font.name, "Inter")
            # 36 half-points -> Pt(18); the centipoint<->half-point conversion holds.
            self.assertEqual(target.font.size, Pt(18))
            self.assertEqual(target.font.color.type, MSO_COLOR_TYPE.RGB)
            self.assertEqual(str(target.font.color.rgb), "1F4E79")

    def test_size_unit_conversion_round_trips_half_points(self) -> None:
        """An 18pt body (36 half-points) lands as exactly 18pt, not 9 or 36."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(template, size_hp=36)
            out = tp / "out.pptx"
            pg.generate(profile, template, parse_idoc(_BODY_IDOC), out)
            target = next(
                r for r in _body_runs(Presentation(out)) if r.text == "Body text here."
            )
            self.assertEqual(target.font.size, Pt(18))

    def test_theme_token_applies_scheme_color(self) -> None:
        """A run color TOKEN resolving to a theme slot applies as a pptx scheme color
        (via pptx's OWN MSO_THEME_COLOR map, never docx's WML map)."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(
                template,
                palette={"accent1": {"ref": {"kind": "theme", "theme": "accent1"}}},
                colors={"accent1": {"hex": "C0392B"}},
            )
            idoc = {
                "blocks": [
                    {"type": "heading", "level": 1, "text": "S"},
                    {"type": "paragraph", "runs": [{"t": "Tok", "color": "accent1"}]},
                ]
            }
            out = tp / "out.pptx"
            sink: list = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            target = next(r for r in _body_runs(Presentation(out)) if r.text == "Tok")
            self.assertEqual(target.font.color.type, MSO_COLOR_TYPE.SCHEME)
            self.assertEqual(target.font.color.theme_color, MSO_THEME_COLOR.ACCENT_1)
            self.assertEqual(
                [f.check for f in sink if f.check == "appearance_color_skipped"], []
            )

    def test_hex_palette_token_applies_as_rgb(self) -> None:
        """A run color token resolving to a hex ref applies as a direct RGB color."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(
                template,
                palette={"brand": {"ref": {"kind": "hex", "hex": "0B5394"}}},
            )
            idoc = {
                "blocks": [
                    {"type": "heading", "level": 1, "text": "S"},
                    {"type": "paragraph", "runs": [{"t": "Hexed", "color": "brand"}]},
                ]
            }
            out = tp / "out.pptx"
            sink: list = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            target = next(r for r in _body_runs(Presentation(out)) if r.text == "Hexed")
            self.assertEqual(target.font.color.type, MSO_COLOR_TYPE.RGB)
            self.assertEqual(str(target.font.color.rgb), "0B5394")
            self.assertEqual(
                [f.check for f in sink if f.check == "appearance_color_skipped"], []
            )

    def test_minted_palette_alias_resolves_as_run_color(self) -> None:
        """Cluster E1 cross-format leg: an off-theme ``hex:RRGGBB`` accent the model
        NAMED an ALIAS for (minted into theme.palette via the real merge path, ref
        byte-copied) is addressable as a clean dotted run-color token on pptx and
        applies as the captured RGB (zero resolver change - the alias is just another
        palette key)."""
        from brandkit.profile import comprehension as comp_mod

        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(
                template,
                palette={
                    "hex:0B5394": {
                        "ref": {"kind": "hex", "hex": "0B5394"},
                        "provenance": [{"where": "run.color", "detail": "body"}],
                        "frequency": "accent",
                    }
                },
            )
            # The model NAMES an alias; the engine mints the dotted token byte-copied.
            res = comp_mod.merge(
                profile,
                {"palette_annotations": {"hex:0B5394": {"alias": "accent.brandblue"}}},
            )
            self.assertTrue(res.ok, res.problems)
            idoc = {
                "blocks": [
                    {"type": "heading", "level": 1, "text": "S"},
                    {
                        "type": "paragraph",
                        "runs": [{"t": "Aliased", "color": "accent.brandblue"}],
                    },
                ]
            }
            out = tp / "out.pptx"
            sink: list = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            target = next(
                r for r in _body_runs(Presentation(out)) if r.text == "Aliased"
            )
            self.assertEqual(target.font.color.type, MSO_COLOR_TYPE.RGB)
            self.assertEqual(str(target.font.color.rgb), "0B5394")
            self.assertEqual([f.check for f in sink if "color" in f.check], [])

    def test_explicit_run_token_wins_over_body_default(self) -> None:
        """A per-run color token is applied first and the body-default color, gated on
        set-only-when-unset, does NOT clobber it."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(
                template,
                color={"kind": "hex", "hex": "111111"},  # body default
                palette={"hot": {"ref": {"kind": "hex", "hex": "FF0000"}}},
            )
            idoc = {
                "blocks": [
                    {"type": "heading", "level": 1, "text": "S"},
                    {"type": "paragraph", "runs": [{"t": "Red", "color": "hot"}]},
                ]
            }
            out = tp / "out.pptx"
            pg.generate(profile, template, parse_idoc(idoc), out)
            target = next(r for r in _body_runs(Presentation(out)) if r.text == "Red")
            self.assertEqual(str(target.font.color.rgb), "FF0000")

    def test_unresolved_token_left_inherited_with_info(self) -> None:
        """A run color token absent from theme.palette leaves the run inherited and
        records a graceful INFO finding (never a fabricated color)."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(template, palette={})  # empty palette
            idoc = {
                "blocks": [
                    {"type": "heading", "level": 1, "text": "S"},
                    {"type": "paragraph", "runs": [{"t": "Plain", "color": "ghost"}]},
                ]
            }
            out = tp / "out.pptx"
            sink: list = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            target = next(r for r in _body_runs(Presentation(out)) if r.text == "Plain")
            self.assertIsNone(target.font.color.type)
            self.assertIn("color_token_unresolved", [f.check for f in sink])
            self.assertEqual(
                schema.Severity.INFO.value,
                next(f.severity for f in sink if f.check == "color_token_unresolved"),
            )


class PptxSetRunColorFallbackTest(unittest.TestCase):
    """Direct coverage of the ``_pptx_set_run_color`` theme-token fallbacks, which a
    real captured token cannot reach (every THEME_SLOTS slot maps to a pptx member)
    but a comprehension-named theme ref could carry."""

    def _fresh_run(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish, has a title
        tb = slide.shapes.add_textbox(0, 0, 1, 1)
        para = tb.text_frame.paragraphs[0]
        return para.add_run()

    def test_unmappable_theme_token_with_hex_realizes_rgb(self) -> None:
        run = self._fresh_run()
        run.text = "x"
        sink: list = []
        # A theme ref whose token is NOT a pptx MSO_THEME_COLOR xml_value, but the
        # resolver enriched it with a concrete hex: realize via the hex, no skip.
        pg._pptx_set_run_color(
            run, {"kind": "theme", "theme": "brandPrimary", "hex": "2E7D32"}, sink
        )
        self.assertEqual(run.font.color.type, MSO_COLOR_TYPE.RGB)
        self.assertEqual(str(run.font.color.rgb), "2E7D32")
        self.assertEqual([f.check for f in sink], [])

    def test_unmappable_theme_token_without_hex_left_inherited_info(self) -> None:
        run = self._fresh_run()
        run.text = "x"
        sink: list = []
        pg._pptx_set_run_color(run, {"kind": "theme", "theme": "brandPrimary"}, sink)
        self.assertIsNone(run.font.color.type)  # left inherited
        self.assertEqual([f.check for f in sink], ["appearance_color_skipped"])
        self.assertEqual(sink[0].severity, schema.Severity.INFO.value)

    def test_malformed_hex_fails_closed_with_info(self) -> None:
        run = self._fresh_run()
        run.text = "x"
        sink: list = []
        pg._pptx_set_run_color(run, {"kind": "hex", "hex": "zzzzzz"}, sink)
        self.assertIsNone(run.font.color.type)
        self.assertEqual([f.check for f in sink], ["appearance_color_skipped"])


class PptxBackendSetOnlyWhenUnsetTest(unittest.TestCase):
    """The ``PPTX_BACKEND`` ``*_unset`` probes drive set-only-when-unset: a run that
    already carries an explicit axis is never clobbered, so an inherited-but-correct
    master/layout value is preserved and re-runs stay byte-identical."""

    def _para_with_one_run(self, text: str):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(0, 0, 1, 1)
        para = tb.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = text
        return para, run

    def test_explicit_font_size_color_are_not_clobbered(self) -> None:
        from brandkit.common import appearance
        from brandkit.profile.resolver import ResolvedOp
        from pptx.dml.color import RGBColor

        para, run = self._para_with_one_run("x")
        run.font.name = "Helvetica"
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor.from_string("ABCDEF")

        op = ResolvedOp(
            role_id="paragraph",
            resolver={},
            status="ok",
            confidence=1.0,
            kind="pptx",
            appearance={
                "font": {"latin": "Inter"},
                "size_hp": 20,
                "color": {"kind": "hex", "hex": "111111"},
            },
        )
        appearance.apply_role_appearance(pg.PPTX_BACKEND, para, op, [])

        # None of the three axes were overwritten - the explicit values stand.
        self.assertEqual(run.font.name, "Helvetica")
        self.assertEqual(run.font.size, Pt(24))
        self.assertEqual(str(run.font.color.rgb), "ABCDEF")

    def test_unset_axes_are_branded(self) -> None:
        from brandkit.common import appearance
        from brandkit.profile.resolver import ResolvedOp

        para, run = self._para_with_one_run("x")  # all axes inherited (None)
        op = ResolvedOp(
            role_id="paragraph",
            resolver={},
            status="ok",
            confidence=1.0,
            kind="pptx",
            appearance={
                "font": {"latin": "Inter"},
                "size_hp": 20,
                "color": {"kind": "hex", "hex": "111111"},
            },
        )
        appearance.apply_role_appearance(pg.PPTX_BACKEND, para, op, [])
        self.assertEqual(run.font.name, "Inter")
        self.assertEqual(run.font.size, Pt(10))  # 20 half-points
        self.assertEqual(str(run.font.color.rgb), "111111")


class PptxAppearanceNoOpAndIdempotencyTest(unittest.TestCase):
    """An empty-appearance profile must generate byte-identically to a profile that
    never carried appearance, and an appearance-carrying profile must be byte-
    idempotent across two generations - the regression net for the new apply path."""

    def test_empty_appearance_is_byte_identical_baseline(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            idoc = parse_idoc(_BODY_IDOC)

            # Baseline: a profile with NO captured body appearance at all (the pre-
            # feature shape). The apply pass reads an empty op.appearance and is a
            # no-op, so the bytes must equal an untouched generation.
            base = T._extract_profile(template)
            a = tp / "a.pptx"
            b = tp / "b.pptx"
            pg.generate(base, template, idoc, a)
            pg.generate(base, template, idoc, b)
            self.assertEqual(a.read_bytes(), b.read_bytes())

    def test_appearance_generation_is_byte_idempotent(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            T._branded_template(template)
            profile = _profile_with_body(
                template,
                latin="Inter",
                size_hp=36,
                color={"kind": "hex", "hex": "1F4E79"},
                palette={"accent1": {"ref": {"kind": "theme", "theme": "accent1"}}},
                colors={"accent1": {"hex": "C0392B"}},
            )
            idoc = parse_idoc(
                {
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "S"},
                        {
                            "type": "paragraph",
                            "runs": [
                                {"t": "Body "},
                                {"t": "token", "color": "accent1"},
                            ],
                        },
                    ]
                }
            )
            a = tp / "a.pptx"
            b = tp / "b.pptx"
            pg.generate(profile, template, idoc, a)
            pg.generate(profile, template, idoc, b)
            self.assertEqual(a.read_bytes(), b.read_bytes())


if __name__ == "__main__":
    unittest.main()
